"""
my-travel 사업방향서 PPT 생성 스크립트
PRD(travel-web-system) 및 코드베이스 분석 기반
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ----- 컬러 팔레트 -----
PRIMARY = RGBColor(0xE6, 0x3F, 0x5C)      # 한국 관광 핑크
PRIMARY_DARK = RGBColor(0xB8, 0x2C, 0x46)
ACCENT = RGBColor(0xFF, 0x8C, 0x42)       # 오렌지
BG_LIGHT = RGBColor(0xFF, 0xF5, 0xF7)
BG_PANEL = RGBColor(0xFA, 0xFA, 0xFA)
INK = RGBColor(0x1F, 0x2A, 0x37)
INK_MUTE = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE5, 0xE7, 0xEB)
SUCCESS = RGBColor(0x10, 0xB9, 0x81)
WARN = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Pretendard"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill=WHITE, line=None, line_w=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill=WHITE, line=None, line_w=0.75, radius=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def header_bar(slide, title, subtitle=None, idx=None, total=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=PRIMARY)
    add_text(slide, Inches(0.6), Inches(0.32), Inches(10), Inches(0.6),
             title, size=26, bold=True, color=INK)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.78), Inches(10), Inches(0.4),
                 subtitle, size=13, color=INK_MUTE)
    if idx is not None and total is not None:
        add_text(slide, Inches(11.8), Inches(0.4), Inches(1.3), Inches(0.4),
                 f"{idx:02d} / {total:02d}", size=11, color=INK_MUTE,
                 align=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0.6), Inches(1.25), Inches(12.1), Emu(9525), fill=LINE)


def footer(slide, label):
    add_text(slide, Inches(0.6), Inches(7.15), Inches(8), Inches(0.3),
             "my-travel · 사업방향서 2026", size=9, color=INK_MUTE)
    add_text(slide, Inches(11.6), Inches(7.15), Inches(1.5), Inches(0.3),
             label, size=9, color=INK_MUTE, align=PP_ALIGN.RIGHT)


TOTAL = 18

# ============================================================
# 1. 커버
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)

# 좌측 색면
add_rect(s, 0, 0, Inches(5.2), SLIDE_H, fill=PRIMARY)
add_text(s, Inches(0.6), Inches(0.6), Inches(4), Inches(0.5),
         "my-travel", size=18, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(6.7), Inches(4), Inches(0.4),
         "2026 · Business Direction", size=11, color=WHITE)

# 메인 카피
add_text(s, Inches(5.8), Inches(1.8), Inches(7), Inches(0.6),
         "BUSINESS DIRECTION", size=12, bold=True, color=PRIMARY)
add_text(s, Inches(5.8), Inches(2.3), Inches(7), Inches(1.6),
         "서울을 넘어,\n한국의 모든 지역을 잇다.",
         size=36, bold=True, color=INK)
add_rect(s, Inches(5.8), Inches(4.2), Inches(0.6), Emu(38100), fill=PRIMARY)
add_text(s, Inches(5.8), Inches(4.4), Inches(7), Inches(1.4),
         "내국인과 방한 외국인을 동시에 연결하는\n"
         "다국어 지역 관광 상품 예약 플랫폼",
         size=14, color=INK_MUTE)
add_text(s, Inches(5.8), Inches(6.4), Inches(7), Inches(0.4),
         "Domestic + Inbound · Multi-language · K-Local Experience",
         size=11, bold=True, color=ACCENT)

# ============================================================
# 2. 목차
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "목차", "Table of Contents", 2, TOTAL)

toc = [
    ("01", "Executive Summary", "한 줄 요약과 핵심 가치"),
    ("02", "문제 정의", "왜 지금, 왜 우리인가"),
    ("03", "시장 기회", "방한 관광 시장과 지역 분산"),
    ("04", "솔루션 개요", "플랫폼이 제공하는 가치"),
    ("05", "타깃 사용자", "4-Sided Marketplace"),
    ("06", "비즈니스 모델", "수익 구조와 단가 가정"),
    ("07", "제품 아키텍처", "Laravel + Vue + Multi-locale"),
    ("08", "현재 구현 현황", "코드베이스 기준 진척도"),
    ("09", "차별화 포인트", "경쟁사 대비 우위"),
    ("10", "GTM 전략 / KPI / 로드맵", "실행 계획과 지표"),
]
y0 = Inches(1.55)
for i, (num, title, desc) in enumerate(toc):
    row_y = y0 + Inches(0.5) * i
    add_text(s, Inches(0.8), row_y, Inches(0.6), Inches(0.45),
             num, size=18, bold=True, color=PRIMARY)
    add_text(s, Inches(1.6), row_y + Inches(0.04), Inches(4.5), Inches(0.4),
             title, size=15, bold=True, color=INK)
    add_text(s, Inches(6.5), row_y + Inches(0.07), Inches(6.5), Inches(0.4),
             desc, size=12, color=INK_MUTE)
footer(s, "Contents")

# ============================================================
# 3. Executive Summary
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "01 Executive Summary", "내·외국인을 동시에 잇는 다국어 지역 관광 예약 플랫폼", 3, TOTAL)

# 좌측 큰 카피
add_text(s, Inches(0.6), Inches(1.6), Inches(7.5), Inches(2.2),
         "관광 수요의 60% 이상이\n서울·제주에 몰린다 — 나머지를 키운다.",
         size=24, bold=True, color=INK)
add_text(s, Inches(0.6), Inches(3.4), Inches(7.5), Inches(2.4),
         "my-travel은 한국인 국내여행자와 방한 외국인이\n"
         "같은 사이트에서 자기 언어로 지역 관광 상품을\n"
         "검색·예약·후기까지 처리하는\n"
         "내·외국인 통합 B2B2C 마켓플레이스입니다.",
         size=14, color=INK_MUTE)

# 우측 핵심 지표 카드 3개
cards = [
    ("4-Sided", "관광객 · 사업자\n가이드 · 관리자", PRIMARY),
    ("4 Lang", "한국어 · 영어\n중국어 · 일본어", ACCENT),
    ("3 Types", "당일 · 숙박\n단일 체험", SUCCESS),
]
cx = Inches(8.3)
cy = Inches(1.6)
cw = Inches(1.45)
ch = Inches(1.8)
gap = Inches(0.1)
for i, (big, desc, col) in enumerate(cards):
    x = cx + (cw + gap) * i
    add_round(s, x, cy, cw, ch, fill=col)
    add_text(s, x, cy + Inches(0.35), cw, Inches(0.6),
             big, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, cy + Inches(0.95), cw, Inches(0.8),
             desc, size=10, color=WHITE, align=PP_ALIGN.CENTER)

# 우측 하단: 한 줄 요약 박스
add_round(s, Inches(8.3), Inches(3.7), Inches(4.4), Inches(2.5),
          fill=BG_PANEL, line=LINE)
add_text(s, Inches(8.55), Inches(3.85), Inches(4), Inches(0.4),
         "ONE-LINER", size=10, bold=True, color=PRIMARY)
add_text(s, Inches(8.55), Inches(4.25), Inches(4), Inches(1.9),
         "내국인과 외국인이 자기 언어로 검색하고,\n"
         "지역 사업자가 직접 등록·운영하며,\n"
         "가이드가 현장에서 QR로 체크인하는\n"
         "내·외국인 통합 로컬 투어 플랫폼.",
         size=12, color=INK)

footer(s, "01 Summary")

# ============================================================
# 4. 문제 정의
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "02 문제 정의", "관광객 · 사업자 · 지역경제 — 세 축의 손실", 4, TOTAL)

problems = [
    ("관광 집중화",
     "내·외국인 관광 수요가\n서울·제주에 과도하게 집중",
     "지역의 우수 상품도 발견되지 못한 채 사장"),
    ("정보 접근성",
     "내국인은 흩어진 블로그·카톡으로 예약,\n외국인은 한국어 장벽으로 접근 불가",
     "내·외국인 통합 다국어 채널 부재"),
    ("사업자 도구 부재",
     "내국인 + 외국인 동시에 다루는\n다국어 예약·운영 도구가 없음",
     "내국 채널 / 외국 채널 이중 운영"),
]
cx0 = Inches(0.6)
cy0 = Inches(1.55)
cw0 = Inches(4.0)
ch0 = Inches(4.5)
gap0 = Inches(0.15)
for i, (t, d, b) in enumerate(problems):
    x = cx0 + (cw0 + gap0) * i
    add_round(s, x, cy0, cw0, ch0, fill=WHITE, line=LINE)
    # 번호
    add_round(s, x + Inches(0.3), cy0 + Inches(0.3),
              Inches(0.7), Inches(0.7), fill=PRIMARY, radius=0.5)
    add_text(s, x + Inches(0.3), cy0 + Inches(0.35),
             Inches(0.7), Inches(0.6),
             f"0{i+1}", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), cy0 + Inches(1.2),
             cw0 - Inches(0.6), Inches(0.5),
             t, size=18, bold=True, color=INK)
    add_text(s, x + Inches(0.3), cy0 + Inches(1.85),
             cw0 - Inches(0.6), Inches(1.5),
             d, size=12, color=INK_MUTE)
    add_rect(s, x + Inches(0.3), cy0 + Inches(3.5),
             cw0 - Inches(0.6), Emu(9525), fill=LINE)
    add_text(s, x + Inches(0.3), cy0 + Inches(3.7),
             cw0 - Inches(0.6), Inches(0.6),
             "→ " + b, size=11, bold=True, color=PRIMARY)

add_round(s, Inches(0.6), Inches(6.25), Inches(12.1), Inches(0.7),
          fill=BG_LIGHT, line=PRIMARY)
add_text(s, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.4),
         "WHY NOW · 국내여행 수요 정착 + 방한 외국인 회복 + 지역균형발전 정책 — 세 흐름이 동시에 열린 시점",
         size=12, bold=True, color=PRIMARY)
footer(s, "02 Problem")

# ============================================================
# 5. 시장 기회
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "03 시장 기회", "방한 관광 회복과 지역 분산 흐름", 5, TOTAL)

# 좌측: 핵심 트렌드
add_text(s, Inches(0.6), Inches(1.5), Inches(6), Inches(0.4),
         "MARKET DRIVERS", size=11, bold=True, color=PRIMARY)
trends = [
    ("국내여행 정착", "주말·생활권 단거리 여행 수요 안정화"),
    ("방한 외국인 회복", "K-컬처 · 면세 · 비자 완화 트렌드"),
    ("로컬 체험 수요", "내·외국인 모두 '진짜 로컬'로 이동"),
    ("지역균형발전", "지자체·관광공사 지역 콘텐츠 투자 확대"),
]
ty = Inches(1.95)
for i, (t, d) in enumerate(trends):
    y = ty + Inches(0.95) * i
    add_round(s, Inches(0.6), y, Inches(0.4), Inches(0.4),
              fill=PRIMARY, radius=0.5)
    add_text(s, Inches(0.6), y + Inches(0.02),
             Inches(0.4), Inches(0.4),
             str(i+1), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.15), y, Inches(5.4), Inches(0.4),
             t, size=14, bold=True, color=INK)
    add_text(s, Inches(1.15), y + Inches(0.4), Inches(5.4), Inches(0.5),
             d, size=11, color=INK_MUTE)

# 우측: TAM/SAM/SOM 시각화 (concentric)
panel_x = Inches(7.0)
panel_y = Inches(1.5)
add_round(s, panel_x, panel_y, Inches(5.7), Inches(5.2),
          fill=BG_PANEL, line=LINE)
add_text(s, panel_x + Inches(0.3), panel_y + Inches(0.2),
         Inches(5), Inches(0.4),
         "TAM · SAM · SOM (가정)", size=11, bold=True, color=PRIMARY)

# 큰 원/중간 원/작은 원
cx = panel_x + Inches(2.85)
cy = panel_y + Inches(3.05)
for r_in, col, label, sub in [
    (2.0, RGBColor(0xFE, 0xE2, 0xE2), "TAM", "한국 국내여행 + 방한 외국인 전체"),
    (1.4, RGBColor(0xFC, 0xA5, 0xA5), "SAM", "지역 로컬투어 잠재 수요 (내·외국)"),
    (0.85, PRIMARY, "SOM", "다국어 예약 가능 1차 타깃"),
]:
    add_round(s,
              cx - Inches(r_in), cy - Inches(r_in),
              Inches(r_in * 2), Inches(r_in * 2),
              fill=col, radius=0.5)
    add_text(s, cx - Inches(r_in), cy - Inches(0.15),
             Inches(r_in * 2), Inches(0.4),
             label, size=13, bold=True,
             color=WHITE if r_in == 0.85 else INK, align=PP_ALIGN.CENTER)

# 범례
add_text(s, panel_x + Inches(0.3), panel_y + Inches(4.5),
         Inches(5.2), Inches(0.4),
         "초기 SOM = 국내 2030 + 영어/일/중화권 FIT 외국인 + 지역 체험",
         size=10, color=INK_MUTE, align=PP_ALIGN.CENTER)

footer(s, "03 Market")

# ============================================================
# 6. 솔루션 개요
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "04 솔루션 개요", "탐색 → 예약 → 운영 → 후기까지 한 자리에서", 6, TOTAL)

# 가로 4단 플로우
steps = [
    ("탐색·발견", "다국어 검색 / 지역·테마 필터\n실시간 재고·평점 노출"),
    ("예약·확정", "자동확정 · 승인필요 분기\n날짜/인원/옵션 선택"),
    ("운영·체크인", "QR 체크인 / 일정 캘린더\n가이드-사업자 실시간 알림"),
    ("후기·재방문", "별점·사진 후기 / 위시리스트\n사업자 답변 + 신고/관리"),
]
fx = Inches(0.6)
fy = Inches(1.55)
fw = Inches(2.95)
fh = Inches(2.6)
gapf = Inches(0.15)
for i, (t, d) in enumerate(steps):
    x = fx + (fw + gapf) * i
    add_round(s, x, fy, fw, fh, fill=WHITE, line=LINE)
    add_round(s, x + Inches(0.25), fy + Inches(0.25),
              Inches(0.55), Inches(0.55), fill=PRIMARY, radius=0.5)
    add_text(s, x + Inches(0.25), fy + Inches(0.28),
             Inches(0.55), Inches(0.5),
             str(i+1), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.95), fy + Inches(0.3),
             fw - Inches(1), Inches(0.5),
             t, size=14, bold=True, color=INK)
    add_text(s, x + Inches(0.25), fy + Inches(1.05),
             fw - Inches(0.5), Inches(1.5),
             d, size=11, color=INK_MUTE)
    # 화살표
    if i < 3:
        ax = x + fw + Inches(0.01)
        add_text(s, ax, fy + Inches(1.05), Inches(0.13), Inches(0.4),
                 "▶", size=12, color=PRIMARY)

# 하단: 운영 가치
add_round(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.55),
          fill=BG_LIGHT, line=None)
add_text(s, Inches(0.85), Inches(4.55), Inches(11), Inches(0.4),
         "PLATFORM PROMISES", size=11, bold=True, color=PRIMARY)

promises = [
    ("내·외국 동시 운영", "한국어 1급 + 영·중·일 — 한 사이트 · 한 운영팀"),
    ("실시간 재고", "날짜별 잔여 수량 / 자동확정·승인 분기 라우팅"),
    ("노쇼 관리", "누적 노쇼 카운트 / 임계치 이상 예약 제한 / 관리자 해제"),
    ("역할별 UX", "관광객/사업자/가이드/관리자 4개 대시보드 분리"),
]
py = Inches(5.05)
pw = Inches(2.95)
ph = Inches(1.8)
for i, (t, d) in enumerate(promises):
    x = Inches(0.6) + (pw + Inches(0.05)) * i
    add_text(s, x + Inches(0.25), py, pw - Inches(0.3), Inches(0.4),
             t, size=13, bold=True, color=INK)
    add_text(s, x + Inches(0.25), py + Inches(0.4),
             pw - Inches(0.3), Inches(1.4),
             d, size=11, color=INK_MUTE)

footer(s, "04 Solution")

# ============================================================
# 7. 타깃 사용자 (4-sided)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "05 타깃 사용자", "4-Sided Marketplace", 7, TOTAL)

quads = [
    ("관광객 (Traveler)",
     "내국인 국내여행자 + 방한 외국인 FIT",
     ["다국어 상품 검색·예약", "위시리스트 / 후기 작성", "마이페이지·예약 관리"],
     PRIMARY),
    ("제공자 (Vendor)",
     "지역 관광·체험 사업자",
     ["상품 등록(다국어)", "캘린더 일정·재고 관리", "예약 승인 / 메시지 대응"],
     ACCENT),
    ("가이드 (Guide)",
     "현장 진행 로컬 가이드",
     ["배정 일정 캘린더 조회", "QR 체크인 / 시작·종료", "노쇼 보고"],
     SUCCESS),
    ("관리자 (Admin)",
     "플랫폼 운영팀",
     ["사용자/상품/예약 관리", "신고·노쇼 처리", "통계 대시보드"],
     WARN),
]
qx = Inches(0.6)
qy = Inches(1.55)
qw = Inches(5.95)
qh = Inches(2.55)
gapq = Inches(0.2)
for i, (title, who, items, col) in enumerate(quads):
    r = i // 2
    c = i % 2
    x = qx + (qw + gapq) * c
    y = qy + (qh + gapq) * r
    add_round(s, x, y, qw, qh, fill=WHITE, line=LINE)
    # 좌측 색 바
    add_rect(s, x, y, Inches(0.12), qh, fill=col)
    add_text(s, x + Inches(0.35), y + Inches(0.2),
             qw - Inches(0.5), Inches(0.4),
             title, size=15, bold=True, color=INK)
    add_text(s, x + Inches(0.35), y + Inches(0.65),
             qw - Inches(0.5), Inches(0.4),
             who, size=11, color=INK_MUTE)
    for j, it in enumerate(items):
        add_text(s, x + Inches(0.35), y + Inches(1.1) + Inches(0.42) * j,
                 qw - Inches(0.5), Inches(0.4),
                 "•  " + it, size=11, color=INK)

footer(s, "05 Users")

# ============================================================
# 8. 비즈니스 모델
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "06 비즈니스 모델", "수익 구조와 단가 가정 (v1: 결제 외부 처리)", 8, TOTAL)

# 좌측: 수익 모델
add_text(s, Inches(0.6), Inches(1.5), Inches(6), Inches(0.4),
         "REVENUE STREAMS", size=11, bold=True, color=PRIMARY)
streams = [
    ("커미션 (Commission)", "확정 예약당 거래액의 8~15%", "성과형 · 사업자 진입장벽 낮춤"),
    ("상품 노출 광고", "추천·기획전·검색 상단 노출", "월 정액 / 입찰형"),
    ("프리미엄 사업자", "다국어 자동번역, 정산 대시보드", "월 구독 (Tier)"),
    ("부가서비스", "사진 촬영, 카피라이팅, 가이드 매칭", "건당 또는 패키지"),
]
sy = Inches(1.95)
for i, (t, a, b) in enumerate(streams):
    y = sy + Inches(1.15) * i
    add_round(s, Inches(0.6), y, Inches(6.0), Inches(1.0), fill=BG_PANEL, line=LINE)
    add_round(s, Inches(0.78), y + Inches(0.3),
              Inches(0.4), Inches(0.4), fill=PRIMARY, radius=0.5)
    add_text(s, Inches(0.78), y + Inches(0.32),
             Inches(0.4), Inches(0.4),
             str(i+1), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.3), y + Inches(0.12),
             Inches(5.0), Inches(0.4),
             t, size=13, bold=True, color=INK)
    add_text(s, Inches(1.3), y + Inches(0.5),
             Inches(5.0), Inches(0.3),
             a, size=11, color=INK)
    add_text(s, Inches(1.3), y + Inches(0.78),
             Inches(5.0), Inches(0.3),
             b, size=10, color=INK_MUTE)

# 우측: Unit Economics 가정
add_round(s, Inches(7.0), Inches(1.5), Inches(5.7), Inches(5.4),
          fill=WHITE, line=LINE)
add_text(s, Inches(7.25), Inches(1.65), Inches(5.2), Inches(0.4),
         "UNIT ECONOMICS · 가정 (Year 1)", size=11, bold=True, color=PRIMARY)

ue_rows = [
    ("평균 객단가 (AOV)", "₩ 80,000"),
    ("커미션율", "10%"),
    ("건당 매출 (Take rate)", "₩ 8,000"),
    ("월 예약 건수 목표", "1,500건"),
    ("월 GMV", "₩ 120,000,000"),
    ("월 플랫폼 매출", "₩ 12,000,000"),
    ("Active Vendors", "150 ~ 250"),
    ("MAU (관광객)", "30,000"),
]
ry = Inches(2.15)
for i, (k, v) in enumerate(ue_rows):
    y = ry + Inches(0.5) * i
    bg = BG_PANEL if i % 2 == 0 else WHITE
    add_rect(s, Inches(7.25), y, Inches(5.2), Inches(0.46), fill=bg)
    add_text(s, Inches(7.4), y + Inches(0.08),
             Inches(3.0), Inches(0.3),
             k, size=11, color=INK)
    add_text(s, Inches(10.0), y + Inches(0.08),
             Inches(2.3), Inches(0.3),
             v, size=11, bold=True, color=PRIMARY, align=PP_ALIGN.RIGHT)

add_text(s, Inches(7.25), Inches(6.5), Inches(5.2), Inches(0.3),
         "※ v1은 현장결제 / 외부결제 — 정산 risk 분리",
         size=9, color=INK_MUTE)

footer(s, "06 Business Model")

# ============================================================
# 9. 제품 아키텍처
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "07 제품 아키텍처", "Laravel 12 + Vue 3 Islands · 다국어 라우팅", 9, TOTAL)

# 상단: 스택 칩
chips = [
    ("Backend", "Laravel 12 / PHP 8.2", PRIMARY),
    ("Frontend", "Blade + Vue 3 + Pinia", ACCENT),
    ("Style", "Tailwind v4 + Vite", SUCCESS),
    ("DB / Cache", "MySQL 8 + Redis", WARN),
    ("Auth", "Sanctum (stateful) + Socialite", INK),
]
cy0 = Inches(1.5)
cx0 = Inches(0.6)
cw0 = Inches(2.4)
for i, (k, v, col) in enumerate(chips):
    x = cx0 + (cw0 + Inches(0.05)) * i
    add_round(s, x, cy0, cw0, Inches(0.85), fill=col)
    add_text(s, x, cy0 + Inches(0.1), cw0, Inches(0.3),
             k, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, cy0 + Inches(0.42), cw0, Inches(0.4),
             v, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 중앙: 3-tier 다이어그램
diag_y = Inches(2.65)
# Frontend
add_round(s, Inches(0.6), diag_y, Inches(12.1), Inches(1.05),
          fill=BG_LIGHT, line=LINE)
add_text(s, Inches(0.85), diag_y + Inches(0.1), Inches(3), Inches(0.3),
         "Presentation · Blade Layouts + Vue 3 Islands", size=10, bold=True, color=PRIMARY)
roles = ["Traveler Web (/{locale})", "Vendor Dashboard (/vendor)",
         "Guide Console (/guide)", "Admin (/admin)"]
for i, r in enumerate(roles):
    add_round(s, Inches(0.85) + Inches(2.95) * i, diag_y + Inches(0.45),
              Inches(2.85), Inches(0.5), fill=WHITE, line=LINE)
    add_text(s, Inches(0.85) + Inches(2.95) * i, diag_y + Inches(0.55),
             Inches(2.85), Inches(0.4),
             r, size=10, color=INK, align=PP_ALIGN.CENTER)

# 화살표
add_text(s, Inches(6.5), diag_y + Inches(1.1), Inches(0.3), Inches(0.3),
         "▼", size=12, color=INK_MUTE, align=PP_ALIGN.CENTER)

# Application layer
app_y = diag_y + Inches(1.45)
add_round(s, Inches(0.6), app_y, Inches(12.1), Inches(1.05),
          fill=WHITE, line=LINE)
add_text(s, Inches(0.85), app_y + Inches(0.1), Inches(6), Inches(0.3),
         "Application · Routes (web.php / api.php) + Services + Policies",
         size=10, bold=True, color=PRIMARY)
mods = ["Auth/RBAC", "Product/Translation", "Booking/Schedule",
        "Review/Message", "Notification"]
mw = Inches(2.32)
for i, m in enumerate(mods):
    add_round(s, Inches(0.85) + (mw + Inches(0.03)) * i,
              app_y + Inches(0.45), mw, Inches(0.5), fill=BG_PANEL, line=LINE)
    add_text(s, Inches(0.85) + (mw + Inches(0.03)) * i,
             app_y + Inches(0.55), mw, Inches(0.4),
             m, size=10, color=INK, align=PP_ALIGN.CENTER)

add_text(s, Inches(6.5), app_y + Inches(1.1), Inches(0.3), Inches(0.3),
         "▼", size=12, color=INK_MUTE, align=PP_ALIGN.CENTER)

# Data layer
data_y = app_y + Inches(1.45)
add_round(s, Inches(0.6), data_y, Inches(12.1), Inches(0.95),
          fill=PRIMARY)
add_text(s, Inches(0.85), data_y + Inches(0.12), Inches(6), Inches(0.3),
         "Data · MySQL 8 (Eloquent + *_translations) · Redis · S3 (planned)",
         size=10, bold=True, color=WHITE)
data_cells = ["users / roles", "products + translations",
              "bookings + schedules", "reviews / messages",
              "notifications / reports"]
dw = Inches(2.32)
for i, c in enumerate(data_cells):
    add_round(s, Inches(0.85) + (dw + Inches(0.03)) * i,
              data_y + Inches(0.45), dw, Inches(0.4),
              fill=PRIMARY_DARK)
    add_text(s, Inches(0.85) + (dw + Inches(0.03)) * i,
             data_y + Inches(0.5), dw, Inches(0.4),
             c, size=10, color=WHITE, align=PP_ALIGN.CENTER)

footer(s, "07 Architecture")

# ============================================================
# 10. 현재 구현 현황
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "08 현재 구현 현황", "코드베이스 기준 v1 진척도", 10, TOTAL)

modules = [
    ("인증 / RBAC", 90, "이메일·소셜 로그인, 4개 역할, user.active 미들웨어"),
    ("상품 관리 (다국어)", 85, "Product+Translation, 카테고리, 지역, 이미지/가격/스케줄"),
    ("예약 시스템", 85, "자동확정/승인필요, 취소, 노쇼 카운트, 가이드 배정"),
    ("리뷰 / 메시지", 80, "별점·사진 후기, 사업자 답변, 예약별 메시지 스레드"),
    ("Vendor 콘솔", 75, "상품 CRUD, 일정·재고, 예약 승인, 리뷰 응답"),
    ("Guide 콘솔", 65, "캘린더 일정, QR 체크인, 시작·종료·노쇼"),
    ("Admin 콘솔", 70, "사용자/상품/예약/신고/노쇼/지역/카테고리"),
    ("블로그 / 콘텐츠", 70, "포스트·카테고리·시리즈·태그·댓글·좋아요"),
    ("결제 / 정산", 0, "v1 제외 — v2 도입 예정"),
    ("푸시 / 알림톡", 30, "이메일 알림 골격, 카카오 알림톡 미구현"),
]

# 헤더
add_text(s, Inches(0.6), Inches(1.5), Inches(4), Inches(0.4),
         "모듈", size=11, bold=True, color=INK_MUTE)
add_text(s, Inches(4.7), Inches(1.5), Inches(4), Inches(0.4),
         "진척도", size=11, bold=True, color=INK_MUTE)
add_text(s, Inches(8.8), Inches(1.5), Inches(4), Inches(0.4),
         "비고", size=11, bold=True, color=INK_MUTE)

row_y = Inches(1.95)
for i, (m, p, n) in enumerate(modules):
    y = row_y + Inches(0.48) * i
    add_text(s, Inches(0.6), y, Inches(4), Inches(0.4),
             m, size=12, bold=True, color=INK)
    # 진척도 바
    bar_x = Inches(4.7)
    bar_w = Inches(3.8)
    add_rect(s, bar_x, y + Inches(0.12),
             bar_w, Inches(0.18), fill=LINE)
    fill_w = int(bar_w * (p / 100))
    if p >= 70:
        col = SUCCESS
    elif p >= 30:
        col = WARN
    else:
        col = INK_MUTE
    if fill_w > 0:
        add_rect(s, bar_x, y + Inches(0.12), fill_w, Inches(0.18), fill=col)
    add_text(s, bar_x + bar_w + Inches(0.1), y, Inches(0.6), Inches(0.4),
             f"{p}%", size=11, bold=True, color=col)
    add_text(s, Inches(8.8), y, Inches(4.2), Inches(0.4),
             n, size=10, color=INK_MUTE)

add_round(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.5),
          fill=BG_LIGHT)
add_text(s, Inches(0.8), Inches(6.97), Inches(11.7), Inches(0.3),
         "v1 코어 70%↑ 완성 — 내국인 우선 오픈 가능, 외국인 확장은 v2에서 가속",
         size=11, bold=True, color=PRIMARY)
footer(s, "08 Status")

# ============================================================
# 11. 차별화 포인트
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "09 차별화 포인트", "왜 my-travel인가", 11, TOTAL)

diffs = [
    ("DIFFERENTIATOR 1",
     "내·외국 통합 사이트",
     "내국인용·외국인용 분리 운영이 아닌\n"
     "한 운영팀이 한 데이터셋으로 동시 노출."),
    ("DIFFERENTIATOR 2",
     "다국어 '직접 입력' 구조",
     "기계 번역이 아닌 사업자가 4개 언어 직접 작성.\n"
     "*_translations 분리 스키마로 언어별 품질·SEO 확보."),
    ("DIFFERENTIATOR 3",
     "현장 운영 우선 설계",
     "QR 체크인, 노쇼 자동 누적·제한, 예약별 메시지 스레드.\n"
     "OTA가 아닌 '운영툴'로 시작 → vendor lock-in."),
    ("DIFFERENTIATOR 4",
     "지역 큐레이션 가능",
     "Region/Category 정규화 + 다국어 + 기획전 슬롯.\n"
     "관광공사·지자체 협업 기획전이 1급 객체로 가능."),
]
dx = Inches(0.6)
dy = Inches(1.55)
dw = Inches(5.95)
dh = Inches(2.55)
gapd = Inches(0.2)
for i, (h, t, d) in enumerate(diffs):
    r = i // 2
    c = i % 2
    x = dx + (dw + gapd) * c
    y = dy + (dh + gapd) * r
    add_round(s, x, y, dw, dh, fill=WHITE, line=LINE)
    add_text(s, x + Inches(0.35), y + Inches(0.25),
             dw - Inches(0.5), Inches(0.3),
             h, size=10, bold=True, color=PRIMARY)
    add_text(s, x + Inches(0.35), y + Inches(0.6),
             dw - Inches(0.5), Inches(0.5),
             t, size=17, bold=True, color=INK)
    add_text(s, x + Inches(0.35), y + Inches(1.3),
             dw - Inches(0.5), Inches(1.2),
             d, size=11, color=INK_MUTE)
    # 우측 모티프
    add_round(s, x + dw - Inches(0.85), y + Inches(0.25),
              Inches(0.55), Inches(0.55),
              fill=BG_LIGHT, radius=0.5)
    add_text(s, x + dw - Inches(0.85), y + Inches(0.32),
             Inches(0.55), Inches(0.4),
             str(i+1), size=14, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER)
footer(s, "09 Differentiation")

# ============================================================
# 12. 경쟁 환경
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "09+ 경쟁 환경", "글로벌 OTA · 국내 액티비티 vs my-travel", 12, TOTAL)

competitors = [
    ("Klook · KKday\n(외국인 OTA)", "외국인 트래픽\n다국어 결제", "내국인 시장 약함\n지역 사업자 진입장벽"),
    ("야놀자 · 마이리얼트립\n(내국인 OTA)", "내국인 인지도\n결제·정산", "외국인·다국어 약함\n해외 마케팅 한계"),
    ("Airbnb Experiences", "글로벌 브랜드\nUGC", "한국 지역 카테고리 얕음\n현장 운영툴 약함"),
    ("my-travel", "내·외국인 동시 운영\n다국어 직접입력\n현장 운영툴", "신규 진입\n양면 수요 집행 필요"),
]
cw0 = Inches(2.95)
cx0 = Inches(0.6)
cy0 = Inches(1.6)
for i, (name, strong, weak) in enumerate(competitors):
    x = cx0 + (cw0 + Inches(0.1)) * i
    add_round(s, x, cy0, cw0, Inches(5.2),
              fill=PRIMARY if i == 3 else WHITE,
              line=PRIMARY if i == 3 else LINE)
    head_color = WHITE if i == 3 else INK
    add_text(s, x + Inches(0.25), cy0 + Inches(0.3),
             cw0 - Inches(0.5), Inches(0.9),
             name, size=14, bold=True, color=head_color)
    label_color = WHITE if i == 3 else PRIMARY
    body_color = WHITE if i == 3 else INK_MUTE
    add_text(s, x + Inches(0.25), cy0 + Inches(1.6),
             cw0 - Inches(0.5), Inches(0.3),
             "STRENGTHS", size=10, bold=True, color=label_color)
    add_text(s, x + Inches(0.25), cy0 + Inches(1.95),
             cw0 - Inches(0.5), Inches(1.4),
             strong, size=11, color=body_color)
    add_text(s, x + Inches(0.25), cy0 + Inches(3.3),
             cw0 - Inches(0.5), Inches(0.3),
             "WEAKNESSES", size=10, bold=True, color=label_color)
    add_text(s, x + Inches(0.25), cy0 + Inches(3.65),
             cw0 - Inches(0.5), Inches(1.4),
             weak, size=11, color=body_color)

footer(s, "09+ Competition")

# ============================================================
# 13. GTM 전략
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "10 GTM 전략", "공급-수요 양면 시장 진입 순서", 13, TOTAL)

# 좌측: 단계
phases = [
    ("Phase 1 · 공급 확보",
     "0~3개월",
     ["부산·제주·강원 3개 지역 선정",
      "지자체·관광공사 협업 시범 사업자 50",
      "온보딩·다국어 작성 지원"]),
    ("Phase 2 · 내국인 우선 수요",
     "3~6개월",
     ["국내 2030 SNS·블로그 마케팅",
      "주말 단거리 여행 키워드 SEO",
      "지자체 기획전·공모 콜라보"]),
    ("Phase 3 · 외국인 확장",
     "6~12개월",
     ["KTO·공항·호텔 컨시어지 채널",
      "Reddit·Xiaohongshu 외국인 KOL",
      "결제·정산·알림톡 v2 출시"]),
]
px = Inches(0.6)
py = Inches(1.55)
pw = Inches(7.5)
ph = Inches(1.7)
for i, (t, d, items) in enumerate(phases):
    y = py + (ph + Inches(0.15)) * i
    add_round(s, px, y, pw, ph, fill=WHITE, line=LINE)
    add_round(s, px, y, Inches(0.6), ph, fill=PRIMARY, radius=0.05)
    add_text(s, px, y + Inches(0.25), Inches(0.6), Inches(0.4),
             f"P{i+1}", size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, px + Inches(0.8), y + Inches(0.2),
             Inches(3.5), Inches(0.4),
             t, size=14, bold=True, color=INK)
    add_text(s, px + Inches(0.8), y + Inches(0.55),
             Inches(3.5), Inches(0.3),
             d, size=10, bold=True, color=PRIMARY)
    for j, it in enumerate(items):
        add_text(s, px + Inches(4.3), y + Inches(0.2) + Inches(0.4) * j,
                 Inches(3.0), Inches(0.4),
                 "•  " + it, size=11, color=INK_MUTE)

# 우측: 채널 카드
add_round(s, Inches(8.4), Inches(1.55), Inches(4.3), Inches(5.4),
          fill=BG_PANEL, line=LINE)
add_text(s, Inches(8.6), Inches(1.7), Inches(4), Inches(0.4),
         "ACQUISITION CHANNELS", size=11, bold=True, color=PRIMARY)
channels = [
    ("국내 SNS·블로그", "내국인 2030 — 인스타·네이버"),
    ("KTO·지자체 협력", "공식 채널 노출·기획전"),
    ("호텔·게스트하우스", "체크인 시 QR·할인 코드"),
    ("SEO (다국어)", "지역+체험 longtail 키워드"),
    ("외국인 커뮤니티", "Reddit / Xiaohongshu / KOL"),
]
chy = Inches(2.15)
for i, (k, v) in enumerate(channels):
    y = chy + Inches(0.92) * i
    add_round(s, Inches(8.6), y, Inches(0.4), Inches(0.4),
              fill=PRIMARY, radius=0.5)
    add_text(s, Inches(8.6), y + Inches(0.05),
             Inches(0.4), Inches(0.4),
             str(i+1), size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.1), y, Inches(3.4), Inches(0.4),
             k, size=12, bold=True, color=INK)
    add_text(s, Inches(9.1), y + Inches(0.4),
             Inches(3.4), Inches(0.4),
             v, size=10, color=INK_MUTE)

footer(s, "10 GTM")

# ============================================================
# 14. KPI
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "10+ KPI · 핵심 지표", "PRD success criteria 기반", 14, TOTAL)

kpis = [
    ("월 방문자", "10,000", "Google Analytics", PRIMARY),
    ("가입 전환율", "≥ 5%", "방문자 → 회원", ACCENT),
    ("예약 전환율", "≥ 2%", "방문자 → 예약", SUCCESS),
    ("예약 완료율", "≥ 80%", "완료 / 전체 예약", PRIMARY),
    ("노쇼율", "≤ 10%", "노쇼 / 확정 예약", WARN),
    ("리뷰 작성율", "≥ 20%", "리뷰 / 완료 예약", ACCENT),
    ("평균 평점", "≥ 4.0", "전체 리뷰 평균", SUCCESS),
    ("Active Vendors", "150+", "활성 사업자", PRIMARY),
]
kx = Inches(0.6)
ky = Inches(1.55)
kw = Inches(2.95)
kh = Inches(2.55)
for i, (k, v, d, col) in enumerate(kpis):
    r = i // 4
    c = i % 4
    x = kx + (kw + Inches(0.1)) * c
    y = ky + (kh + Inches(0.15)) * r
    add_round(s, x, y, kw, kh, fill=WHITE, line=LINE)
    add_rect(s, x, y, kw, Inches(0.12), fill=col)
    add_text(s, x + Inches(0.25), y + Inches(0.35),
             kw - Inches(0.5), Inches(0.4),
             k, size=12, bold=True, color=INK_MUTE)
    add_text(s, x + Inches(0.25), y + Inches(0.85),
             kw - Inches(0.5), Inches(1.0),
             v, size=36, bold=True, color=col)
    add_text(s, x + Inches(0.25), y + Inches(2.0),
             kw - Inches(0.5), Inches(0.4),
             d, size=10, color=INK_MUTE)
footer(s, "10+ KPI")

# ============================================================
# 15. 로드맵
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "10++ 로드맵", "v1 출시 → v2 결제·정산 → v3 확장", 15, TOTAL)

phases_r = [
    ("v1.0 · MVP",
     "Q1 2026",
     ["관광객/사업자/가이드/관리자 4 콘솔",
      "다국어 상품 · 예약 · 리뷰 · 메시지",
      "한국어 메인 + 영어, 현장결제 운영"],
     PRIMARY),
    ("v1.5 · 내국인 강화",
     "Q2 2026",
     ["국내 결제(토스·카카오페이) 도입",
      "카카오 알림톡 / SMS",
      "지역 기획전 슬롯 / 큐레이션"],
     ACCENT),
    ("v2.0 · 외국인 확장",
     "Q3 2026",
     ["해외 카드·PayPal 결제",
      "환불 자동화 / 정산 대시보드",
      "중국어·일본어 정식 오픈"],
     SUCCESS),
    ("v3.0 · 확장",
     "Q4 2026~",
     ["AI 추천 / 동적 가격",
      "모바일 앱 (관광객·가이드)",
      "보험·교통·숙박 외부 연동"],
     WARN),
]
rx = Inches(0.6)
ry = Inches(1.6)
rw = Inches(2.95)
rh = Inches(5.0)
for i, (t, q, items, col) in enumerate(phases_r):
    x = rx + (rw + Inches(0.1)) * i
    add_round(s, x, ry, rw, rh, fill=WHITE, line=LINE)
    add_rect(s, x, ry, rw, Inches(0.65), fill=col)
    add_text(s, x + Inches(0.25), ry + Inches(0.1),
             rw - Inches(0.5), Inches(0.5),
             t, size=14, bold=True, color=WHITE)
    add_text(s, x + Inches(0.25), ry + Inches(0.85),
             rw - Inches(0.5), Inches(0.4),
             q, size=12, bold=True, color=col)
    for j, it in enumerate(items):
        add_text(s, x + Inches(0.25), ry + Inches(1.4) + Inches(0.9) * j,
                 rw - Inches(0.5), Inches(0.9),
                 "•  " + it, size=11, color=INK)

footer(s, "10++ Roadmap")

# ============================================================
# 16. 리스크
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "리스크 & 대응", "PRD 식별 위험 + 운영 가정", 16, TOTAL)

risks = [
    ("노쇼 증가", "높음", "노쇼 누적 카운트 + 임계치 자동 제한 + 관리자 해제 흐름 구현 완료"),
    ("다국어 품질", "중간", "사업자 직접 입력 + 네이티브 검수 인력 / 자동번역은 v2에서 보조"),
    ("사업자 이탈", "높음", "운영툴(캘린더·체크인·메시지) 우선 → vendor lock-in 효과"),
    ("트래픽 급증", "중간", "Redis 캐시 + S3/CDN(예정) + 수평 확장 설계"),
    ("부적절한 리뷰", "중간", "신고 시스템 + 관리자 검토 + 신고 다발 사용자 자동 플래그"),
    ("결제·환불 분쟁", "높음", "v1 현장결제로 risk 분리, v2 PG 도입 시 환불 정책 자동화"),
]
add_text(s, Inches(0.6), Inches(1.5), Inches(4.5), Inches(0.4),
         "RISK", size=10, bold=True, color=INK_MUTE)
add_text(s, Inches(5.2), Inches(1.5), Inches(1.5), Inches(0.4),
         "IMPACT", size=10, bold=True, color=INK_MUTE)
add_text(s, Inches(6.9), Inches(1.5), Inches(6), Inches(0.4),
         "MITIGATION", size=10, bold=True, color=INK_MUTE)

ry0 = Inches(1.9)
for i, (r, imp, m) in enumerate(risks):
    y = ry0 + Inches(0.75) * i
    bg = BG_PANEL if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.7), fill=bg)
    add_text(s, Inches(0.8), y + Inches(0.2), Inches(4.3), Inches(0.4),
             r, size=12, bold=True, color=INK)
    # 임팩트 칩
    chip_col = {"높음": PRIMARY, "중간": WARN, "낮음": SUCCESS}[imp]
    add_round(s, Inches(5.2), y + Inches(0.18),
              Inches(0.95), Inches(0.35), fill=chip_col, radius=0.5)
    add_text(s, Inches(5.2), y + Inches(0.2),
             Inches(0.95), Inches(0.3),
             imp, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(6.9), y + Inches(0.2),
             Inches(5.7), Inches(0.45),
             m, size=11, color=INK_MUTE)
footer(s, "Risks")

# ============================================================
# 17. ASK
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
header_bar(s, "ASK · 다음 12개월에 필요한 것", "팀 / 자본 / 파트너십", 17, TOTAL)

# 좌측: ASK 카드들
asks = [
    ("팀",
     ["풀스택 1 (Laravel/Vue)",
      "운영·CS 다국어 1",
      "BD 1 (사업자 온보딩)",
      "디자이너 0.5"]),
    ("자본",
     ["12개월 운영비",
      "초기 마케팅 예산",
      "인프라(AWS S3·CDN·CloudFront)",
      "PG 연동 (v2)"]),
    ("파트너십",
     ["KTO / 지역 관광공사",
      "주요 도시 지자체",
      "공항·호텔 컨시어지",
      "해외 OTA 어필리에이트"]),
]
ax0 = Inches(0.6)
ay0 = Inches(1.55)
aw0 = Inches(4.0)
ah0 = Inches(4.8)
for i, (t, items) in enumerate(asks):
    x = ax0 + (aw0 + Inches(0.15)) * i
    add_round(s, x, ay0, aw0, ah0, fill=WHITE, line=LINE)
    add_rect(s, x, ay0, aw0, Inches(0.7), fill=PRIMARY)
    add_text(s, x + Inches(0.3), ay0 + Inches(0.15),
             aw0 - Inches(0.6), Inches(0.5),
             t, size=18, bold=True, color=WHITE)
    for j, it in enumerate(items):
        y = ay0 + Inches(1.0) + Inches(0.7) * j
        add_round(s, x + Inches(0.3), y + Inches(0.1),
                  Inches(0.3), Inches(0.3), fill=PRIMARY, radius=0.5)
        add_text(s, x + Inches(0.3), y + Inches(0.13),
                 Inches(0.3), Inches(0.3),
                 str(j+1), size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.75), y + Inches(0.1),
                 aw0 - Inches(1), Inches(0.4),
                 it, size=12, color=INK)

# 하단 핵심 메시지
add_round(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.5),
          fill=BG_LIGHT)
add_text(s, Inches(0.85), Inches(6.65), Inches(11.7), Inches(0.3),
         "v1 코어 70%↑ 완성. 다음 12개월은 '공급 100 → 300 + 수요 1만 MAU' 양면 시장 점화 단계.",
         size=12, bold=True, color=PRIMARY)
footer(s, "Ask")

# ============================================================
# 18. 클로징
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, INK)

# 우측 색면
add_rect(s, Inches(8.5), 0, Inches(4.83), SLIDE_H, fill=PRIMARY)

# 좌측 메시지
add_text(s, Inches(0.8), Inches(1.0), Inches(7), Inches(0.4),
         "CLOSING", size=11, bold=True, color=PRIMARY)
add_text(s, Inches(0.8), Inches(1.6), Inches(7.5), Inches(2.6),
         "내국인도, 외국인도\n같은 한국을 본다.",
         size=38, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.5), Inches(7.5), Inches(2.5),
         "한국의 모든 지역과 모든 여행자를\n"
         "다국어로 연결하는 단 하나의 운영툴 ·\n"
         "my-travel.",
         size=16, color=RGBColor(0xD1, 0xD5, 0xDB))
add_text(s, Inches(0.8), Inches(6.5), Inches(7), Inches(0.4),
         "Thank you.", size=14, bold=True, color=PRIMARY)

# 우측 핵심 숫자
add_text(s, Inches(8.85), Inches(1.0), Inches(4.4), Inches(0.4),
         "BY THE NUMBERS", size=11, bold=True, color=WHITE)
big = [
    ("4", "사용자 역할"),
    ("4", "지원 언어"),
    ("3", "상품 유형"),
    ("70%+", "v1 진척도"),
]
by = Inches(1.6)
for i, (n, lab) in enumerate(big):
    y = by + Inches(1.25) * i
    add_text(s, Inches(8.85), y, Inches(4.4), Inches(0.85),
             n, size=42, bold=True, color=WHITE)
    add_text(s, Inches(8.85), y + Inches(0.95),
             Inches(4.4), Inches(0.3),
             lab, size=11, color=RGBColor(0xFE, 0xCD, 0xD3))

# 저장
output_path = "/Users/jisungyoo/Dev/my-travel/docs/my-travel-business-direction.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
