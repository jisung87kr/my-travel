"""
my-travel 서비스 소개 PPT
대상: 사용자/사업자/파트너에게 서비스 자체를 소개
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----- 컬러 팔레트 (서비스 소개는 좀 더 밝고 부드럽게) -----
PRIMARY = RGBColor(0xE6, 0x3F, 0x5C)
PRIMARY_DARK = RGBColor(0xB8, 0x2C, 0x46)
PINK_50 = RGBColor(0xFF, 0xF1, 0xF2)
PINK_100 = RGBColor(0xFF, 0xE4, 0xE6)
PINK_200 = RGBColor(0xFE, 0xCD, 0xD3)
ACCENT = RGBColor(0xFF, 0x8C, 0x42)
ACCENT_LIGHT = RGBColor(0xFF, 0xED, 0xD5)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
TEAL = RGBColor(0x06, 0xB6, 0xD4)
TEAL_LIGHT = RGBColor(0xCF, 0xFA, 0xFE)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
PURPLE_LIGHT = RGBColor(0xEA, 0xE2, 0xFD)
GREEN = RGBColor(0x10, 0xB9, 0x81)
GREEN_LIGHT = RGBColor(0xD1, 0xFA, 0xE5)
YELLOW = RGBColor(0xF5, 0x9E, 0x0B)
YELLOW_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
INK = RGBColor(0x1F, 0x2A, 0x37)
INK_MUTE = RGBColor(0x6B, 0x72, 0x80)
INK_LIGHT = RGBColor(0x9C, 0xA3, 0xAF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFA, 0xFA, 0xFA)
LINE = RGBColor(0xE5, 0xE7, 0xEB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background(); bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Pretendard"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill=WHITE, line=None, line_w=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill=WHITE, line=None, line_w=0.75, radius=0.1):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_circle(slide, x, y, d, fill=WHITE, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def page_header(slide, kicker, title, idx, total):
    add_text(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.4),
             kicker, size=12, bold=True, color=PRIMARY)
    add_text(slide, Inches(0.6), Inches(0.75), Inches(11), Inches(0.7),
             title, size=28, bold=True, color=INK)
    add_text(slide, Inches(11.8), Inches(0.5), Inches(1.3), Inches(0.4),
             f"{idx:02d} / {total:02d}", size=11, color=INK_LIGHT,
             align=PP_ALIGN.RIGHT)
    # underline accent
    add_rect(slide, Inches(0.6), Inches(1.55), Inches(0.6), Emu(28575), fill=PRIMARY)


def footer(slide, label=""):
    add_text(slide, Inches(0.6), Inches(7.15), Inches(6), Inches(0.3),
             "my-travel · Service Introduction", size=9, color=INK_LIGHT)
    add_text(slide, Inches(11.6), Inches(7.15), Inches(1.5), Inches(0.3),
             label, size=9, color=INK_LIGHT, align=PP_ALIGN.RIGHT)


TOTAL = 17

# ============================================================
# 1. 커버
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)

# 배경 데코 원들
add_circle(s, Inches(9.5), Inches(-2.5), Inches(7), fill=PINK_50)
add_circle(s, Inches(10.8), Inches(3.5), Inches(4), fill=PINK_100)
add_circle(s, Inches(11.5), Inches(5.8), Inches(2.2), fill=PRIMARY)
add_circle(s, Inches(8.8), Inches(5.2), Inches(1.2), fill=ACCENT)

# 로고/브랜드
add_circle(s, Inches(0.8), Inches(0.7), Inches(0.5), fill=PRIMARY)
add_text(s, Inches(0.8), Inches(0.78), Inches(0.5), Inches(0.4),
         "m", size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.45), Inches(0.85), Inches(4), Inches(0.4),
         "my-travel", size=18, bold=True, color=INK)

add_text(s, Inches(0.8), Inches(7.0), Inches(6), Inches(0.4),
         "SERVICE INTRODUCTION  ·  2026", size=11, color=INK_LIGHT)

# 메인 카피
add_rect(s, Inches(0.8), Inches(2.4), Inches(0.6), Emu(28575), fill=PRIMARY)
add_text(s, Inches(0.8), Inches(2.55), Inches(8), Inches(0.5),
         "HELLO, KOREA · 안녕, 한국",
         size=13, bold=True, color=PRIMARY)

add_text(s, Inches(0.8), Inches(3.0), Inches(9), Inches(1.6),
         "한국을 여행하는 모두를 위한\n지역 관광 예약 플랫폼",
         size=40, bold=True, color=INK)

add_text(s, Inches(0.8), Inches(5.0), Inches(8), Inches(1.2),
         "내국인도, 외국인도 — 같은 사이트에서\n"
         "서울 너머 한국 구석구석을 경험하세요.",
         size=16, color=INK_MUTE)

# 언어 칩
langs = [("KO", "한국어"), ("EN", "English"), ("中", "中文"), ("日", "日本語")]
ly = Inches(6.3)
for i, (k, v) in enumerate(langs):
    x = Inches(0.8) + Inches(1.55) * i
    add_round(s, x, ly, Inches(1.4), Inches(0.5),
              fill=WHITE, line=LINE, radius=0.4)
    add_circle(s, x + Inches(0.12), ly + Inches(0.13),
               Inches(0.24), fill=PRIMARY)
    add_text(s, x + Inches(0.12), ly + Inches(0.18),
             Inches(0.24), Inches(0.2),
             k, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.45), ly + Inches(0.16),
             Inches(0.9), Inches(0.3),
             v, size=10, color=INK)

# ============================================================
# 2. 목차
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "CONTENTS", "이 자료에서 다루는 것", 2, TOTAL)

toc = [
    ("01", "서비스 개요", "my-travel은 어떤 서비스인가요?"),
    ("02", "왜 만들었나요", "어떤 문제를 풀고 싶었는지"),
    ("03", "누구를 위한 서비스", "관광객 · 사업자 · 가이드"),
    ("04", "핵심 가치", "다국어 · 실시간 · 로컬"),
    ("05", "관광객 사용법", "검색 → 예약 → 후기"),
    ("06", "사업자 사용법", "등록 → 일정 → 예약 관리"),
    ("07", "가이드 사용법", "일정 → QR 체크인 → 보고"),
    ("08", "상품 라인업", "투어 · 액티비티 · 체험"),
    ("09", "비즈니스 모델", "어떻게 운영되는가"),
    ("10", "왜 my-travel인가", "다른 곳과 무엇이 다른가"),
]
y0 = Inches(2.0)
col_w = Inches(6.0)
for i, (num, title, desc) in enumerate(toc):
    r = i // 2
    c = i % 2
    y = y0 + Inches(0.5) * r
    x = Inches(0.6) + col_w * c
    add_text(s, x, y, Inches(0.6), Inches(0.4),
             num, size=18, bold=True, color=PRIMARY)
    add_text(s, x + Inches(0.8), y + Inches(0.03),
             Inches(2.8), Inches(0.4),
             title, size=13, bold=True, color=INK)
    add_text(s, x + Inches(0.8), y + Inches(0.32),
             Inches(4.5), Inches(0.25),
             desc, size=10, color=INK_MUTE)
footer(s, "Contents")

# ============================================================
# 3. 서비스 개요 - What is my-travel
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "01  ABOUT", "my-travel은 어떤 서비스인가요?", 3, TOTAL)

# 좌측 일러스트 영역
add_round(s, Inches(0.6), Inches(2.0), Inches(5.5), Inches(4.7),
          fill=PINK_50, radius=0.06)
# 가짜 폰 화면 모티프
phone_x = Inches(2.0)
phone_y = Inches(2.3)
add_round(s, phone_x, phone_y, Inches(2.7), Inches(4.1),
          fill=WHITE, line=LINE, radius=0.08)
add_rect(s, phone_x + Inches(0.15), phone_y + Inches(0.2),
         Inches(2.4), Inches(0.4), fill=PRIMARY)
add_text(s, phone_x + Inches(0.3), phone_y + Inches(0.28),
         Inches(2), Inches(0.3),
         "어디로 떠나시나요?", size=9, bold=True, color=WHITE)
# 카드 더미들
for i in range(3):
    cy = phone_y + Inches(0.75) + Inches(1.05) * i
    add_round(s, phone_x + Inches(0.15), cy,
              Inches(2.4), Inches(0.95), fill=BG, radius=0.08)
    add_rect(s, phone_x + Inches(0.25), cy + Inches(0.1),
             Inches(0.8), Inches(0.5),
             fill=[PINK_200, ACCENT_LIGHT, BLUE_LIGHT][i])
    add_text(s, phone_x + Inches(1.15), cy + Inches(0.12),
             Inches(1.3), Inches(0.3),
             ["서울 야경 투어", "제주 한라산 트레킹", "부산 미식 투어"][i],
             size=8, bold=True, color=INK)
    add_text(s, phone_x + Inches(1.15), cy + Inches(0.4),
             Inches(1.3), Inches(0.3),
             "★ 4.9 · 12,000원~", size=7, color=INK_MUTE)

# 우측 설명
add_text(s, Inches(6.6), Inches(2.0), Inches(6.5), Inches(0.4),
         "ONE-LINE", size=11, bold=True, color=PRIMARY)
add_text(s, Inches(6.6), Inches(2.4), Inches(6.5), Inches(1.8),
         "내국인과 외국인이\n"
         "같은 사이트, 자기 언어로\n"
         "한국 지역 관광을 예약하는 플랫폼",
         size=22, bold=True, color=INK)

add_text(s, Inches(6.6), Inches(4.5), Inches(6.5), Inches(2.0),
         "서울·제주 패키지 일변도를 넘어,\n"
         "부산·강원·경주·전주 등 전국의\n"
         "로컬 사업자가 운영하는 진짜 한국 경험을\n"
         "내·외국인 모두에게 한 자리에서.",
         size=13, color=INK_MUTE)
footer(s, "01 About")

# ============================================================
# 4. 왜 만들었나 - Mission
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "02  WHY", "왜 my-travel을 만들었나요?", 4, TOTAL)

# 큰 인용구
add_round(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.2),
          fill=PINK_50, radius=0.04)
add_text(s, Inches(0.95), Inches(2.1), Inches(0.6), Inches(0.6),
         "“", size=60, bold=True, color=PRIMARY)
add_text(s, Inches(1.85), Inches(2.3), Inches(10.5), Inches(1.0),
         "내국인도 외국인도 — 한국을 더 깊게 보고 싶지만,\n로컬 경험은 여전히 흩어져 있고 찾기 어렵습니다.",
         size=22, bold=True, color=INK)
add_text(s, Inches(1.85), Inches(3.65), Inches(10.5), Inches(0.5),
         "— 우리는 이 거리를 좁히고 싶습니다.",
         size=13, color=INK_MUTE)

# 하단 3개 미션
missions = [
    ("관광객에게는", "발견의 기쁨을",
     "내국인·외국인 모두\n언어 장벽 없이 로컬 발견", PRIMARY),
    ("사업자에게는", "새로운 고객을",
     "내국인 + 외국인을 동시에 만나는\n다국어 디지털 창구", ACCENT),
    ("지역에게는", "관광의 분산을",
     "서울·제주 집중 관광 수익이\n전국으로 흘러가도록", GREEN),
]
my0 = Inches(4.6)
mw = Inches(3.95)
mh = Inches(2.4)
for i, (k, t, d, col) in enumerate(missions):
    x = Inches(0.6) + (mw + Inches(0.15)) * i
    add_round(s, x, my0, mw, mh, fill=WHITE, line=LINE)
    add_circle(s, x + Inches(0.3), my0 + Inches(0.3),
               Inches(0.4), fill=col)
    add_text(s, x + Inches(0.3), my0 + Inches(0.35),
             Inches(0.4), Inches(0.3),
             str(i+1), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.9), my0 + Inches(0.32),
             mw - Inches(1), Inches(0.3),
             k, size=11, color=INK_MUTE)
    add_text(s, x + Inches(0.3), my0 + Inches(0.95),
             mw - Inches(0.6), Inches(0.5),
             t, size=16, bold=True, color=INK)
    add_text(s, x + Inches(0.3), my0 + Inches(1.55),
             mw - Inches(0.6), Inches(0.7),
             d, size=11, color=INK_MUTE)
footer(s, "02 Why")

# ============================================================
# 5. 누구를 위한 서비스
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "03  FOR WHOM", "누구를 위한 서비스인가요?", 5, TOTAL)

personas = [
    ("관광객\nTraveler",
     "내국인 + 방한 외국인 여행자",
     ["• 주말·휴가 국내여행을 찾는 한국인",
      "• 서울 밖을 가보고 싶은 외국인",
      "• 자기 언어로 안심 예약하고 싶은 분"],
     PRIMARY, PINK_50),
    ("사업자\nVendor",
     "지역에서 관광/체험을 운영하는 분",
     ["• 내국인 + 외국인 모두 만나고 싶은 분",
      "• 예약·일정 관리를 디지털화하고 싶은 분",
      "• 카카오톡 수기 관리에서 벗어나고 싶은 분"],
     ACCENT, ACCENT_LIGHT),
    ("로컬 가이드\nGuide",
     "현장에서 투어를 진행하는 분",
     ["• 배정 일정을 캘린더로 확인",
      "• QR 체크인으로 빠르게 출석 확인",
      "• 노쇼·특이사항 바로 보고"],
     GREEN, GREEN_LIGHT),
]
py = Inches(2.0)
pw = Inches(4.0)
ph = Inches(4.7)
for i, (name, who, items, col, bg) in enumerate(personas):
    x = Inches(0.6) + (pw + Inches(0.1)) * i
    add_round(s, x, py, pw, ph, fill=WHITE, line=LINE)
    # 상단 컬러 영역
    add_round(s, x, py, pw, Inches(1.85), fill=bg)
    # 아바타 원
    add_circle(s, x + Inches(1.65), py + Inches(0.4),
               Inches(0.7), fill=col)
    add_text(s, x + Inches(1.65), py + Inches(0.55),
             Inches(0.7), Inches(0.4),
             ["A", "B", "G"][i], size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, x, py + Inches(1.2), pw, Inches(0.6),
             name, size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), py + Inches(2.0),
             pw - Inches(0.6), Inches(0.4),
             who, size=11, bold=True, color=col)
    for j, it in enumerate(items):
        add_text(s, x + Inches(0.3), py + Inches(2.55) + Inches(0.55) * j,
                 pw - Inches(0.6), Inches(0.5),
                 it, size=11, color=INK_MUTE)
footer(s, "03 For Whom")

# ============================================================
# 6. 핵심 가치
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "04  CORE VALUES", "이 서비스가 약속하는 것", 6, TOTAL)

values = [
    ("내·외국인 동시 지원", "ALL-IN-ONE",
     "한국어가 1급 시민\n+ 외국어 동시 운영\n한 사이트 / 한 운영팀",
     "한 · 영 · 중 · 일", PRIMARY, PINK_50),
    ("실시간 재고", "REAL-TIME",
     "날짜별 잔여 수량을\n바로 확인하고 예약\n늦었다 / 매진 걱정 없음",
     "DAY-BY-DAY", BLUE, BLUE_LIGHT),
    ("로컬 우선", "AUTHENTIC LOCAL",
     "여행사 대행이 아니라\n현지에서 직접 운영하는\n진짜 로컬 사업자",
     "BY LOCAL HOSTS", GREEN, GREEN_LIGHT),
    ("현장 안심", "ON-SITE CARE",
     "QR 체크인 · 메시지 · 노쇼 관리\n예약부터 종료까지\n전 과정을 책임",
     "END-TO-END", PURPLE, PURPLE_LIGHT),
]
vx0 = Inches(0.6)
vy0 = Inches(2.0)
vw = Inches(2.95)
vh = Inches(4.65)
for i, (k, en, d, tag, col, bg) in enumerate(values):
    x = vx0 + (vw + Inches(0.1)) * i
    add_round(s, x, vy0, vw, vh, fill=bg)
    # 큰 숫자
    add_text(s, x + Inches(0.3), vy0 + Inches(0.25),
             vw - Inches(0.5), Inches(0.6),
             f"0{i+1}", size=40, bold=True, color=col)
    add_text(s, x + Inches(0.3), vy0 + Inches(1.0),
             vw - Inches(0.5), Inches(0.4),
             en, size=10, bold=True, color=col)
    add_text(s, x + Inches(0.3), vy0 + Inches(1.45),
             vw - Inches(0.5), Inches(0.5),
             k, size=18, bold=True, color=INK)
    add_text(s, x + Inches(0.3), vy0 + Inches(2.25),
             vw - Inches(0.5), Inches(1.5),
             d, size=11, color=INK_MUTE)
    # 하단 태그
    add_round(s, x + Inches(0.3), vy0 + vh - Inches(0.85),
              vw - Inches(0.6), Inches(0.5),
              fill=WHITE, radius=0.4)
    add_text(s, x + Inches(0.3), vy0 + vh - Inches(0.73),
             vw - Inches(0.6), Inches(0.3),
             tag, size=10, bold=True, color=col, align=PP_ALIGN.CENTER)
footer(s, "04 Values")

# ============================================================
# 7. 관광객 사용법
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "05  TRAVELER FLOW", "관광객은 이렇게 사용합니다", 7, TOTAL)

flow = [
    ("검색", "SEARCH",
     "지역 · 날짜 · 인원으로\n원하는 체험 찾기",
     "🔍", PRIMARY, PINK_50),
    ("선택", "BROWSE",
     "사진 · 후기 · 가격\n포함사항까지 확인",
     "👀", ACCENT, ACCENT_LIGHT),
    ("예약", "BOOK",
     "옵션 선택 후 예약\n즉시확정 또는 승인",
     "📅", BLUE, BLUE_LIGHT),
    ("체험", "EXPERIENCE",
     "현장에서 QR로 체크인\n가이드와 함께 출발",
     "🎒", GREEN, GREEN_LIGHT),
    ("공유", "REVIEW",
     "별점 · 사진 후기로\n다음 여행자에게 추천",
     "⭐", PURPLE, PURPLE_LIGHT),
]
fy = Inches(2.5)
fw = Inches(2.35)
fh = Inches(3.6)
gap = Inches(0.05)
for i, (k, en, d, emo, col, bg) in enumerate(flow):
    x = Inches(0.6) + (fw + gap) * i
    add_round(s, x, fy, fw, fh, fill=bg)
    add_circle(s, x + Inches(0.85), fy + Inches(0.3),
               Inches(0.65), fill=col)
    add_text(s, x + Inches(0.85), fy + Inches(0.42),
             Inches(0.65), Inches(0.4),
             str(i+1), size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), fy + Inches(1.15),
             fw - Inches(0.4), Inches(0.4),
             en, size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), fy + Inches(1.5),
             fw - Inches(0.4), Inches(0.5),
             k, size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), fy + Inches(2.25),
             fw - Inches(0.4), Inches(1.2),
             d, size=11, color=INK_MUTE, align=PP_ALIGN.CENTER)

# 화살표 라인
add_text(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.5),
         "전 과정을 모국어로 · 모바일 한 화면에서 완결",
         size=13, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
footer(s, "05 Traveler")

# ============================================================
# 8. 사업자 사용법
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "06  VENDOR FLOW", "사업자는 이렇게 운영합니다", 8, TOTAL)

# 좌측: 4단계 카드
v_steps = [
    ("상품 등록",
     "사진 · 가격 · 옵션 · 다국어 설명\n자동확정 또는 승인필요 선택",
     ACCENT),
    ("일정 관리",
     "캘린더에서 날짜별 수량 설정\n특정 날짜 운영 중단 가능",
     BLUE),
    ("예약 처리",
     "신청 도착 시 알림\n승인 / 거절 / 메시지 전송",
     GREEN),
    ("후기 응답",
     "고객 리뷰 확인 및 답변\n부적절 후기 신고 처리",
     PURPLE),
]
sy = Inches(2.0)
sx = Inches(0.6)
sw = Inches(6.5)
sh = Inches(1.15)
for i, (t, d, col) in enumerate(v_steps):
    y = sy + (sh + Inches(0.12)) * i
    add_round(s, sx, y, sw, sh, fill=WHITE, line=LINE)
    add_rect(s, sx, y, Inches(0.15), sh, fill=col)
    add_round(s, sx + Inches(0.35), y + Inches(0.35),
              Inches(0.55), Inches(0.55), fill=col, radius=0.5)
    add_text(s, sx + Inches(0.35), y + Inches(0.4),
             Inches(0.55), Inches(0.5),
             str(i+1), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, sx + Inches(1.1), y + Inches(0.15),
             sw - Inches(1.3), Inches(0.4),
             t, size=15, bold=True, color=INK)
    add_text(s, sx + Inches(1.1), y + Inches(0.55),
             sw - Inches(1.3), Inches(0.7),
             d, size=11, color=INK_MUTE)

# 우측: 사업자가 얻는 가치
add_round(s, Inches(7.4), Inches(2.0), Inches(5.3), Inches(4.85),
          fill=ACCENT_LIGHT, radius=0.04)
add_text(s, Inches(7.6), Inches(2.15), Inches(5), Inches(0.4),
         "WHAT YOU GET", size=11, bold=True, color=ACCENT)
add_text(s, Inches(7.6), Inches(2.55), Inches(5), Inches(0.8),
         "사업자에게 제공되는 것",
         size=20, bold=True, color=INK)

benefits = [
    "내국인 + 외국인 고객 동시 접점",
    "한·영·중·일 4개 언어 상품 노출",
    "실시간 캘린더 · 재고 관리",
    "QR 기반 현장 체크인 도구",
    "예약별 메시지 스레드",
    "노쇼 자동 카운트 · 예약 제한",
    "리뷰 응답 / 관리자 분쟁 중재",
]
by0 = Inches(3.65)
for i, b in enumerate(benefits):
    y = by0 + Inches(0.42) * i
    add_circle(s, Inches(7.6), y + Inches(0.1),
               Inches(0.22), fill=ACCENT)
    add_text(s, Inches(7.6), y + Inches(0.12),
             Inches(0.22), Inches(0.2),
             "✓", size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.95), y + Inches(0.08),
             Inches(4.5), Inches(0.4),
             b, size=12, color=INK)
footer(s, "06 Vendor")

# ============================================================
# 9. 가이드 사용법
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "07  GUIDE FLOW", "가이드는 이렇게 진행합니다", 9, TOTAL)

# 좌측: 워크플로우
gx = Inches(0.6)
gy = Inches(2.0)
add_round(s, gx, gy, Inches(7.0), Inches(4.85),
          fill=GREEN_LIGHT, radius=0.04)
add_text(s, gx + Inches(0.3), gy + Inches(0.2),
         Inches(6.5), Inches(0.4),
         "GUIDE DAY", size=11, bold=True, color=GREEN)
add_text(s, gx + Inches(0.3), gy + Inches(0.6),
         Inches(6.5), Inches(0.5),
         "투어 당일의 표준 플로우",
         size=20, bold=True, color=INK)

g_steps = [
    ("아침", "오늘 배정 투어 확인", "캘린더에서 일정 · 고객 명단 확인"),
    ("집결", "고객 QR 체크인", "QR 스캔으로 참석자 출석 처리"),
    ("진행", "투어 시작 처리", "시작 버튼으로 상태 업데이트"),
    ("종료", "투어 완료 보고", "완료 / 노쇼 / 특이사항 메모"),
]
for i, (when, t, d) in enumerate(g_steps):
    y = gy + Inches(1.4) + Inches(0.85) * i
    add_round(s, gx + Inches(0.4), y + Inches(0.1),
              Inches(0.7), Inches(0.65), fill=GREEN, radius=0.5)
    add_text(s, gx + Inches(0.4), y + Inches(0.2),
             Inches(0.7), Inches(0.4),
             when, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, gx + Inches(0.4), y + Inches(0.4),
             Inches(0.7), Inches(0.3),
             "▼", size=8, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, gx + Inches(1.3), y + Inches(0.1),
             Inches(5.5), Inches(0.4),
             t, size=14, bold=True, color=INK)
    add_text(s, gx + Inches(1.3), y + Inches(0.45),
             Inches(5.5), Inches(0.4),
             d, size=11, color=INK_MUTE)

# 우측: QR 콘셉트 일러스트
add_round(s, Inches(7.8), Inches(2.0), Inches(4.9), Inches(4.85),
          fill=WHITE, line=LINE)
add_text(s, Inches(8.0), Inches(2.15), Inches(4.6), Inches(0.4),
         "QR CHECK-IN", size=11, bold=True, color=PRIMARY)
add_text(s, Inches(8.0), Inches(2.55), Inches(4.6), Inches(0.5),
         "스캔 한 번으로 출석 완료",
         size=18, bold=True, color=INK)

# QR 미니어처 (격자)
qr_x = Inches(8.6)
qr_y = Inches(3.45)
qr_d = Inches(2.6)
add_round(s, qr_x, qr_y, qr_d, qr_d, fill=WHITE, line=INK, radius=0.06)
import random
random.seed(42)
cell = Inches(0.2)
inner_x = qr_x + Inches(0.3)
inner_y = qr_y + Inches(0.3)
for i in range(10):
    for j in range(10):
        if random.random() > 0.55:
            add_rect(s, inner_x + cell * j, inner_y + cell * i,
                     cell, cell, fill=INK)
# 코너 마커
for cx_, cy_ in [(inner_x, inner_y),
                  (inner_x + cell * 8, inner_y),
                  (inner_x, inner_y + cell * 8)]:
    add_rect(s, cx_, cy_, cell * 2, cell * 2, fill=WHITE)
    add_rect(s, cx_, cy_, cell * 2, cell * 2, fill=INK)
    add_rect(s, cx_ + cell * 0.3, cy_ + cell * 0.3,
             cell * 1.4, cell * 1.4, fill=WHITE)
    add_rect(s, cx_ + cell * 0.6, cy_ + cell * 0.6,
             cell * 0.8, cell * 0.8, fill=INK)

add_text(s, Inches(8.0), Inches(6.3), Inches(4.6), Inches(0.4),
         "고객의 예약 확인 화면을 스캔하면\n자동으로 체크인 + 시작 처리",
         size=11, color=INK_MUTE, align=PP_ALIGN.CENTER)
footer(s, "07 Guide")

# ============================================================
# 10. 상품 라인업
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "08  PRODUCT LINEUP", "어떤 상품을 만날 수 있나요?", 10, TOTAL)

# 상단: 3가지 상품 유형
types = [
    ("당일 코스",
     "One-Day Tour",
     "하루 안에 끝나는\n패키지형 투어",
     "예) 서울 야경 + 남산 + 한강",
     PRIMARY),
    ("숙박 패키지",
     "Multi-Day Package",
     "1박 이상의\n체류형 상품",
     "예) 제주 2박 3일 힐링 스테이",
     BLUE),
    ("단일 체험",
     "Single Activity",
     "짧고 캐주얼한\n원포인트 체험",
     "예) 한복 체험 · 쿠킹 클래스",
     ACCENT),
]
tx = Inches(0.6)
ty = Inches(2.0)
tw = Inches(4.0)
th = Inches(2.0)
for i, (k, en, d, ex, col) in enumerate(types):
    x = tx + (tw + Inches(0.15)) * i
    add_round(s, x, ty, tw, th, fill=WHITE, line=LINE)
    add_rect(s, x, ty, Inches(0.15), th, fill=col)
    add_text(s, x + Inches(0.35), ty + Inches(0.25),
             tw - Inches(0.5), Inches(0.3),
             en, size=10, bold=True, color=col)
    add_text(s, x + Inches(0.35), ty + Inches(0.6),
             tw - Inches(0.5), Inches(0.5),
             k, size=18, bold=True, color=INK)
    add_text(s, x + Inches(0.35), ty + Inches(1.1),
             tw - Inches(0.5), Inches(0.6),
             d, size=11, color=INK_MUTE)
    add_text(s, x + Inches(0.35), ty + Inches(1.6),
             tw - Inches(0.5), Inches(0.35),
             ex, size=10, color=col, bold=True)

# 하단: 카테고리/지역
cy0 = Inches(4.4)
add_round(s, Inches(0.6), cy0, Inches(6.05), Inches(2.5),
          fill=PINK_50, radius=0.04)
add_text(s, Inches(0.85), cy0 + Inches(0.2),
         Inches(5.5), Inches(0.4),
         "카테고리", size=11, bold=True, color=PRIMARY)
cats = [("투어", PRIMARY), ("액티비티", ACCENT), ("문화체험", BLUE),
        ("식도락", GREEN), ("자연", TEAL), ("야간투어", PURPLE),
        ("티켓", YELLOW), ("패키지", PRIMARY_DARK)]
for i, (c, col) in enumerate(cats):
    r = i // 4
    cc = i % 4
    x = Inches(0.85) + Inches(1.35) * cc
    y = cy0 + Inches(0.65) + Inches(0.7) * r
    add_round(s, x, y, Inches(1.25), Inches(0.5),
              fill=WHITE, line=col, radius=0.4)
    add_text(s, x, y + Inches(0.12), Inches(1.25), Inches(0.3),
             c, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)

add_round(s, Inches(6.75), cy0, Inches(5.95), Inches(2.5),
          fill=BLUE_LIGHT, radius=0.04)
add_text(s, Inches(7.0), cy0 + Inches(0.2),
         Inches(5.5), Inches(0.4),
         "지역", size=11, bold=True, color=BLUE)
regs = ["서울", "부산", "제주", "경기", "강원", "경주", "전주", "기타"]
for i, r in enumerate(regs):
    rr = i // 4
    cc = i % 4
    x = Inches(7.0) + Inches(1.35) * cc
    y = cy0 + Inches(0.65) + Inches(0.7) * rr
    add_round(s, x, y, Inches(1.25), Inches(0.5),
              fill=WHITE, line=BLUE, radius=0.4)
    add_text(s, x, y + Inches(0.12), Inches(1.25), Inches(0.3),
             r, size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
footer(s, "08 Lineup")

# ============================================================
# 11. 비즈니스 모델
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "09  BUSINESS MODEL", "어떻게 운영되나요?", 11, TOTAL)

# 좌측: 흐름도
add_text(s, Inches(0.6), Inches(2.0), Inches(8), Inches(0.4),
         "MARKETPLACE FLOW", size=11, bold=True, color=PRIMARY)
add_text(s, Inches(0.6), Inches(2.4), Inches(8), Inches(0.6),
         "관광객 ↔ 사업자를 연결하는 양면 시장",
         size=18, bold=True, color=INK)

# 3-노드 다이어그램
node_y = Inches(3.8)
nodes = [
    ("관광객", "TRAVELER", PRIMARY, PINK_50),
    ("my-travel", "PLATFORM", INK, BG),
    ("사업자", "VENDOR", ACCENT, ACCENT_LIGHT),
]
nx0 = Inches(0.6)
nw = Inches(2.4)
for i, (k, en, col, bg) in enumerate(nodes):
    x = nx0 + (nw + Inches(0.4)) * i
    add_round(s, x, node_y, nw, Inches(1.5), fill=bg, line=col)
    add_text(s, x, node_y + Inches(0.3), nw, Inches(0.4),
             en, size=10, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, x, node_y + Inches(0.65), nw, Inches(0.6),
             k, size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    if i < 2:
        ax = x + nw
        add_text(s, ax, node_y + Inches(0.55),
                 Inches(0.4), Inches(0.4),
                 "↔", size=20, bold=True, color=INK_LIGHT,
                 align=PP_ALIGN.CENTER)

# 흐름 설명
add_text(s, Inches(0.6), Inches(5.7), Inches(8), Inches(0.4),
         "예약이 확정되면 → 거래액의 일부가 my-travel의 운영 비용",
         size=12, color=INK_MUTE)
add_text(s, Inches(0.6), Inches(6.1), Inches(8), Inches(0.4),
         "v1은 현장 결제 / v2부터 온라인 결제 도입 예정",
         size=11, color=INK_LIGHT)

# 우측: 수익 카드들
add_round(s, Inches(8.7), Inches(2.0), Inches(4.0), Inches(4.85),
          fill=WHITE, line=LINE)
add_text(s, Inches(8.9), Inches(2.15), Inches(3.7), Inches(0.4),
         "REVENUE STREAMS", size=11, bold=True, color=PRIMARY)
add_text(s, Inches(8.9), Inches(2.55), Inches(3.7), Inches(0.5),
         "수익 구조",
         size=18, bold=True, color=INK)

rev = [
    ("커미션", "확정 예약의\n8~15%", PRIMARY),
    ("노출 광고", "추천·기획전\n상단 노출", ACCENT),
    ("프리미엄", "월 구독\n부가기능 제공", BLUE),
    ("부가서비스", "촬영·번역\n패키지", GREEN),
]
for i, (k, d, col) in enumerate(rev):
    r = i // 2
    c = i % 2
    x = Inches(8.9) + Inches(1.85) * c
    y = Inches(3.55) + Inches(1.55) * r
    add_round(s, x, y, Inches(1.7), Inches(1.4), fill=BG, line=LINE)
    add_circle(s, x + Inches(0.25), y + Inches(0.2),
               Inches(0.35), fill=col)
    add_text(s, x + Inches(0.25), y + Inches(0.27),
             Inches(0.35), Inches(0.3),
             str(i+1), size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.7), y + Inches(0.2),
             Inches(1.0), Inches(0.4),
             k, size=12, bold=True, color=INK)
    add_text(s, x + Inches(0.2), y + Inches(0.75),
             Inches(1.5), Inches(0.6),
             d, size=10, color=INK_MUTE, align=PP_ALIGN.CENTER)
footer(s, "09 Model")

# ============================================================
# 12. 왜 my-travel인가 - 차별화
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "10  WHY US", "다른 서비스와 무엇이 다른가요?", 12, TOTAL)

# 비교 카드
items = [
    ("일반 OTA / 여행사", "my-travel",
     [
         ("내국인용·외국인용 사이트 분리", "내·외국인 한 사이트 동시 운영"),
         ("기계 번역 / 한국어만", "사업자 직접 4개 언어 작성"),
         ("대형 패키지 중심", "지역 로컬 사업자 중심"),
         ("수기 예약 / 카톡 문의", "실시간 캘린더 + 메시지"),
         ("노쇼 시 손실 일방 부담", "노쇼 자동 카운트 + 제한 정책"),
     ]),
]
# 헤더
add_round(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(0.7),
          fill=BG, line=LINE)
add_text(s, Inches(0.85), Inches(2.18), Inches(5.5), Inches(0.4),
         items[0][0], size=14, bold=True, color=INK_LIGHT, align=PP_ALIGN.CENTER)

add_round(s, Inches(6.65), Inches(2.0), Inches(6.05), Inches(0.7),
          fill=PRIMARY)
add_text(s, Inches(6.9), Inches(2.18), Inches(5.5), Inches(0.4),
         items[0][1], size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 비교 행
ry0 = Inches(2.85)
rh = Inches(0.78)
for i, (a, b) in enumerate(items[0][2]):
    y = ry0 + (rh + Inches(0.05)) * i
    add_round(s, Inches(0.6), y, Inches(5.8), rh, fill=WHITE, line=LINE)
    add_text(s, Inches(0.95), y + Inches(0.25),
             Inches(5.4), Inches(0.5),
             "✕ " + a, size=12, color=INK_MUTE)
    add_round(s, Inches(6.65), y, Inches(6.05), rh,
              fill=PINK_50, line=PRIMARY)
    add_text(s, Inches(7.0), y + Inches(0.25),
             Inches(5.7), Inches(0.5),
             "✓ " + b, size=12, bold=True, color=PRIMARY_DARK)
footer(s, "10 Why Us")

# ============================================================
# 13. 기술 & 안정성
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "11  RELIABILITY", "안정적인 기술 위에 만들었습니다", 13, TOTAL)

tech = [
    ("검증된 스택",
     "Laravel 12 + Vue 3 + MySQL\n전 세계 수만 서비스가 사용 중",
     "PROVEN", PRIMARY),
    ("빠른 응답",
     "Redis 캐시 + 페이지 < 3초\nAPI < 500ms 목표",
     "FAST", BLUE),
    ("안전한 인증",
     "Sanctum · OAuth (Google/Kakao/Apple)\nCSRF · bcrypt · HTTPS 필수",
     "SECURE", GREEN),
    ("확장 가능",
     "S3/CDN · 수평 확장 설계\n트래픽 증가에 대응 가능",
     "SCALABLE", PURPLE),
    ("다국어 SEO",
     "*_translations 분리 스키마\n언어별 검색 최적화",
     "MULTI-LANG", ACCENT),
    ("운영 도구",
     "관리자 통계 · 신고 · 노쇼\n분쟁 중재 흐름 내장",
     "OPS-READY", TEAL),
]
ty = Inches(2.0)
tw = Inches(4.0)
th = Inches(2.35)
for i, (k, d, tag, col) in enumerate(tech):
    r = i // 3
    c = i % 3
    x = Inches(0.6) + (tw + Inches(0.1)) * c
    y = ty + (th + Inches(0.15)) * r
    add_round(s, x, y, tw, th, fill=WHITE, line=LINE)
    add_round(s, x + Inches(0.3), y + Inches(0.3),
              Inches(1.2), Inches(0.35), fill=col, radius=0.4)
    add_text(s, x + Inches(0.3), y + Inches(0.36),
             Inches(1.2), Inches(0.3),
             tag, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), y + Inches(0.85),
             tw - Inches(0.5), Inches(0.4),
             k, size=15, bold=True, color=INK)
    add_text(s, x + Inches(0.3), y + Inches(1.35),
             tw - Inches(0.5), Inches(0.9),
             d, size=11, color=INK_MUTE)
footer(s, "11 Reliability")

# ============================================================
# 14. 다국어 지원
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "12  MULTI-LANGUAGE", "한 사이트, 네 가지 언어", 14, TOTAL)

# 좌측: 4언어 카드들
langs = [
    ("한국어", "Korean", "전국 구석구석을 더 가까이",
     "내국인 메인", "KO", PRIMARY, PINK_50, True),
    ("English", "English", "Discover Korea, beyond Seoul.",
     "방한 외국인", "EN", BLUE, BLUE_LIGHT, True),
    ("中文", "Chinese", "发现首尔之外的韩国魅力",
     "중화권", "中", ACCENT, ACCENT_LIGHT, False),
    ("日本語", "Japanese", "ソウルを越えた韓国の魅力",
     "일본", "日", PURPLE, PURPLE_LIGHT, False),
]
ly0 = Inches(2.0)
lw = Inches(3.0)
lh = Inches(4.85)
for i, (k, en, cp, region, code, col, bg, active) in enumerate(langs):
    x = Inches(0.6) + (lw + Inches(0.1)) * i
    add_round(s, x, ly0, lw, lh, fill=bg)
    add_circle(s, x + Inches(1.05), ly0 + Inches(0.4),
               Inches(0.9), fill=col)
    add_text(s, x + Inches(1.05), ly0 + Inches(0.55),
             Inches(0.9), Inches(0.5),
             code, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, ly0 + Inches(1.5), lw, Inches(0.5),
             k, size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, x, ly0 + Inches(2.05), lw, Inches(0.3),
             en, size=10, color=INK_MUTE, align=PP_ALIGN.CENTER)
    add_rect(s, x + Inches(1.1), ly0 + Inches(2.55),
             Inches(0.8), Emu(9525), fill=col)
    add_text(s, x + Inches(0.3), ly0 + Inches(2.85),
             lw - Inches(0.6), Inches(1.5),
             cp, size=12, color=INK, align=PP_ALIGN.CENTER)
    # 상태 칩
    chip_col = col if active else INK_LIGHT
    chip_text = "활성" if active else "준비중"
    add_round(s, x + Inches(0.9), ly0 + lh - Inches(0.7),
              Inches(1.2), Inches(0.4),
              fill=WHITE, line=chip_col, radius=0.5)
    add_text(s, x + Inches(0.9), ly0 + lh - Inches(0.6),
             Inches(1.2), Inches(0.3),
             chip_text, size=10, bold=True, color=chip_col, align=PP_ALIGN.CENTER)
footer(s, "12 Languages")

# ============================================================
# 15. 안전 장치
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "13  TRUST & SAFETY", "안심하고 쓸 수 있도록", 15, TOTAL)

safety = [
    ("신원 확인",
     "사업자 등록 정보 · 관리자 승인을 거친\n사업자만 상품을 등록",
     "🛡", PRIMARY, PINK_50),
    ("실시간 메시지",
     "예약별 메시지 스레드로\n문의·변동사항을 즉시 공유",
     "💬", BLUE, BLUE_LIGHT),
    ("리뷰 신고",
     "부적절한 후기는 신고 후\n관리자가 검토하여 조치",
     "⚠", ACCENT, ACCENT_LIGHT),
    ("노쇼 관리",
     "누적 노쇼 시 예약 제한\n사업자도 동일한 정책으로 보호",
     "🚫", PURPLE, PURPLE_LIGHT),
    ("관리자 중재",
     "분쟁 발생 시 관리자가\n중립적으로 예약·환불 처리",
     "⚖", GREEN, GREEN_LIGHT),
    ("개인정보 보호",
     "HTTPS · bcrypt · CSRF\n표준 보안 수칙 모두 적용",
     "🔒", TEAL, TEAL_LIGHT),
]
sy = Inches(2.0)
sw = Inches(4.0)
sh = Inches(2.35)
for i, (k, d, emo, col, bg) in enumerate(safety):
    r = i // 3
    c = i % 3
    x = Inches(0.6) + (sw + Inches(0.1)) * c
    y = sy + (sh + Inches(0.15)) * r
    add_round(s, x, y, sw, sh, fill=bg)
    add_circle(s, x + Inches(0.3), y + Inches(0.3),
               Inches(0.55), fill=col)
    add_text(s, x + Inches(0.3), y + Inches(0.42),
             Inches(0.55), Inches(0.4),
             str(i+1), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), y + Inches(1.0),
             sw - Inches(0.5), Inches(0.4),
             k, size=15, bold=True, color=INK)
    add_text(s, x + Inches(0.3), y + Inches(1.45),
             sw - Inches(0.5), Inches(0.9),
             d, size=11, color=INK_MUTE)
footer(s, "13 Safety")

# ============================================================
# 16. 향후 계획
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
page_header(s, "14  WHAT'S NEXT", "앞으로 이렇게 발전합니다", 16, TOTAL)

# 가로 타임라인
phases = [
    ("NOW",
     "v1 출시",
     "관광객/사업자/가이드/관리자 4 콘솔\n다국어 · 예약 · 리뷰 · 메시지",
     PRIMARY),
    ("SOON",
     "운영 강화",
     "카카오 알림톡 · 지역 기획전\n다국어 작성 지원",
     ACCENT),
    ("NEXT",
     "결제 · 정산",
     "카드 · 해외결제\n환불 자동화 · 정산 대시보드",
     BLUE),
    ("FUTURE",
     "확장",
     "AI 추천 · 모바일 앱\n외부 연동(보험·교통·숙박)",
     PURPLE),
]
tw = Inches(2.95)
ty = Inches(2.2)
th = Inches(4.0)
for i, (when, t, d, col) in enumerate(phases):
    x = Inches(0.6) + (tw + Inches(0.1)) * i
    add_round(s, x, ty, tw, th, fill=WHITE, line=LINE)
    # 상단 칩
    add_round(s, x + Inches(0.3), ty + Inches(0.3),
              Inches(1.2), Inches(0.4), fill=col, radius=0.5)
    add_text(s, x + Inches(0.3), ty + Inches(0.4),
             Inches(1.2), Inches(0.3),
             when, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 제목
    add_text(s, x + Inches(0.3), ty + Inches(1.0),
             tw - Inches(0.5), Inches(0.6),
             t, size=18, bold=True, color=INK)
    # 설명
    add_text(s, x + Inches(0.3), ty + Inches(1.8),
             tw - Inches(0.5), Inches(2.0),
             d, size=12, color=INK_MUTE)
    # 큰 화살표
    if i < 3:
        ax = x + tw + Inches(0.01)
        add_text(s, ax, ty + Inches(1.9),
                 Inches(0.13), Inches(0.4),
                 "▶", size=14, color=PRIMARY)

# 하단 메시지
add_round(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.55),
          fill=PINK_50)
add_text(s, Inches(0.85), Inches(6.62), Inches(11.7), Inches(0.3),
         "지금은 시작 · 1년 안에 결제까지 · 그 다음은 한국을 넘어 동아시아로",
         size=13, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
footer(s, "14 Roadmap")

# ============================================================
# 17. 클로징
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)

# 좌측 색면
add_rect(s, 0, 0, Inches(6.5), SH, fill=PRIMARY)
# 데코
add_circle(s, Inches(-1.5), Inches(5.5), Inches(4), fill=PRIMARY_DARK)
add_circle(s, Inches(4.5), Inches(-1), Inches(3.5), fill=ACCENT)

add_text(s, Inches(0.7), Inches(0.8), Inches(5), Inches(0.4),
         "my-travel", size=16, bold=True, color=WHITE)
add_rect(s, Inches(0.7), Inches(2.2), Inches(0.6), Emu(28575), fill=WHITE)
add_text(s, Inches(0.7), Inches(2.35), Inches(5.5), Inches(0.4),
         "JOIN US", size=12, bold=True, color=PINK_200)
add_text(s, Inches(0.7), Inches(2.8), Inches(5.5), Inches(2.5),
         "함께\n한국을\n다시 그립니다.",
         size=44, bold=True, color=WHITE)

add_text(s, Inches(0.7), Inches(6.5), Inches(5), Inches(0.4),
         "Hello, my-travel.", size=14, bold=True, color=PINK_100)

# 우측 정보
add_text(s, Inches(7.0), Inches(0.8), Inches(5.5), Inches(0.4),
         "TRY IT NOW", size=11, bold=True, color=PRIMARY)
add_text(s, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.8),
         "지금 시작하세요",
         size=30, bold=True, color=INK)

# CTA 박스들
ctas = [
    ("관광객으로 시작",
     "내국인·외국인 모두 회원가입 후\n원하는 지역의 상품을 검색해보세요",
     PRIMARY, "FOR TRAVELER"),
    ("사업자로 등록",
     "사업자 등록 후 첫 상품을\n등록해보세요",
     ACCENT, "FOR VENDOR"),
    ("가이드로 합류",
     "지역 운영팀에게\n가이드 신청을 보내주세요",
     GREEN, "FOR GUIDE"),
]
cy = Inches(2.7)
for i, (t, d, col, tag) in enumerate(ctas):
    y = cy + Inches(1.35) * i
    add_round(s, Inches(7.0), y, Inches(5.7), Inches(1.2),
              fill=WHITE, line=LINE)
    add_rect(s, Inches(7.0), y, Inches(0.15), Inches(1.2), fill=col)
    add_text(s, Inches(7.3), y + Inches(0.15),
             Inches(5), Inches(0.3),
             tag, size=9, bold=True, color=col)
    add_text(s, Inches(7.3), y + Inches(0.45),
             Inches(5), Inches(0.4),
             t, size=14, bold=True, color=INK)
    add_text(s, Inches(7.3), y + Inches(0.8),
             Inches(5), Inches(0.4),
             d, size=10, color=INK_MUTE)

add_text(s, Inches(7.0), Inches(6.85), Inches(5.5), Inches(0.3),
         "Thank you.  ·  Made with ♥ in Korea",
         size=11, color=INK_LIGHT)

# 저장
output_path = "/Users/jisungyoo/Dev/my-travel/docs/my-travel-service-intro.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
